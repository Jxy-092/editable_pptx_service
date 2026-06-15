from io import BytesIO
from pptx import Presentation


class PPTXBuilder:
    def __init__(self):
        self.prs = None

    def create_presentation(self):
        self.prs = Presentation()

    def setup_presentation_size(self, width_pixels: int, height_pixels: int):
        ratio = width_pixels / max(1, height_pixels)
        self.prs.slide_width = int(9144000)
        self.prs.slide_height = int(9144000 / ratio)

    def add_blank_slide(self):
        blank = self.prs.slide_layouts[6]
        return self.prs.slides.add_slide(blank)

    def save(self, output_file: str):
        self.prs.save(output_file)

    def to_bytes(self) -> bytes:
        bio = BytesIO()
        self.prs.save(bio)
        return bio.getvalue()
